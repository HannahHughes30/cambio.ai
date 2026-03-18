"""Build the final presentation for the Cambio AI Agent project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Slide dimensions (standard 10x5.63 widescreen)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.63)

# Color palette
BG = RGBColor(0x1A, 0x1A, 0x2E)       # dark navy
ACCENT = RGBColor(0x00, 0xD2, 0xFF)    # cyan
ACCENT2 = RGBColor(0xFF, 0x6B, 0x6B)   # coral
ACCENT3 = RGBColor(0x4E, 0xCB, 0x71)   # green
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xCC, 0xCC, 0xCC)     # light gray
DIM = RGBColor(0x88, 0x88, 0xAA)       # dim text
CARD_BG = RGBColor(0x2A, 0x2A, 0x44)   # card background
DARK_CARD = RGBColor(0x22, 0x22, 0x3A)


def add_bg(slide):
    """Add dark background to slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    color=LIGHT, bold_items=None, sub_items=None):
    """Add a bulleted list. bold_items is a set of indices to bold. sub_items maps index -> list of sub-bullets."""
    if bold_items is None:
        bold_items = set()
    if sub_items is None:
        sub_items = {}

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = ACCENT if i in bold_items else color
        p.font.bold = i in bold_items
        p.font.name = 'Calibri'
        p.space_after = Pt(4)
        p.level = 0

        if i in sub_items:
            for sub in sub_items[i]:
                sp = tf.add_paragraph()
                sp.text = sub
                sp.font.size = Pt(font_size - 2)
                sp.font.color.rgb = DIM
                sp.font.name = 'Calibri'
                sp.space_after = Pt(2)
                sp.level = 1

    return tf


def add_card(slide, left, top, width, height, title, body_lines, accent_color=ACCENT):
    """Add a rounded rectangle card with title and body."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = accent_color
    p.font.bold = True
    p.font.name = 'Calibri'
    p.space_after = Pt(6)

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT
        p.font.name = 'Calibri'
        p.space_after = Pt(2)


def add_section_header(slide, number, title):
    """Add a section number + title at the top."""
    add_text_box(slide, Inches(0.5), Inches(0.3), Inches(1), Inches(0.5),
                 f"0{number}" if number < 10 else str(number),
                 font_size=36, color=ACCENT, bold=True)
    add_text_box(slide, Inches(0.5), Inches(0.75), Inches(8), Inches(0.6),
                 title, font_size=28, color=WHITE, bold=True)
    # Accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(2), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_stat_box(slide, left, top, value, label, color=ACCENT):
    """Add a big stat number with label."""
    add_text_box(slide, left, top, Inches(1.8), Inches(0.6),
                 value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, top + Inches(0.55), Inches(1.8), Inches(0.4),
                 label, font_size=10, color=DIM, alignment=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = Emu(9144000)   # 10 inches
    prs.slide_height = Emu(5143500)  # 5.63 inches
    blank_layout = prs.slide_layouts[6]  # blank

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(8.4), Inches(1.2),
                 "Design and Evaluation of\nAI Agents for Cambio",
                 font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.8), Inches(2.6), Inches(8), Inches(0.5),
                 "Belief-State Tracking, Expected Value Reasoning, and Opponent Modeling\nin a Partially Observable Card Game",
                 font_size=16, color=ACCENT)
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(8), Inches(0.4),
                 "Nas Ahmed  ·  Hannah Hughes  ·  Juan E. Cisneros",
                 font_size=14, color=LIGHT)
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(8), Inches(0.4),
                 "California Polytechnic State University, San Luis Obispo",
                 font_size=12, color=DIM)
    # Accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.45), Inches(3), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    # =========================================================================
    # SLIDE 2: The Challenge
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 1, "The Challenge")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8),
                 "How far can classical probabilistic reasoning go\nin imperfect-information games?",
                 font_size=20, color=ACCENT, bold=True)

    items = [
        "Most cards are face-down → agents can't see the full game state",
        "Stochastic card draws → can't plan perfectly even with perfect memory",
        "Multi-agent competition → each player's strategy affects everyone",
        "Memory-based gameplay → must track every observation across turns",
        "Strategic timing → calling Cambio too early or too late is costly",
        "We started with RL and MCTS in mind. We didn't need them."
    ]
    add_bullet_list(slide, Inches(0.7), Inches(2.5), Inches(8.5), Inches(2.8),
                    items, font_size=14, bold_items={5})

    # =========================================================================
    # SLIDE 3: What is Cambio?
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 2, "What is Cambio?")

    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.2), Inches(3.6),
             "CORE RULES", [
                 "• 2-6 players, 54-card deck (52 + 2 Jokers)",
                 "• Each player starts with 4 face-down cards",
                 "• Peek at your bottom 2 cards at start",
                 "• Goal: minimize total hand value",
                 "",
                 "• Each turn: draw from deck or discard,",
                 "  then swap into hand OR discard",
                 "• Discard pile draw → must swap",
                 "",
                 "• Call 'Cambio' → everyone gets 1 last turn",
                 "• Caller bonus: -10 if lowest, +10 if not",
             ])

    add_card(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(1.7),
             "POWER CARDS (triggered on discard)", [
                 "7/8 → Peek at one of YOUR cards",
                 "9/10 → Peek at an OPPONENT's card",
                 "J/Q → Blind swap any two cards",
                 "Black King → Peek ANY card + optional swap",
             ], accent_color=ACCENT2)

    add_card(slide, Inches(5.0), Inches(3.4), Inches(4.5), Inches(1.7),
             "SPECIAL VALUES", [
                 "Red King (Hearts/Diamonds) = -1 point",
                 "Joker = 0 points  ·  Ace = 1 point",
                 "J/Q/Black King = 10 points each",
                 "Number cards = face value",
             ], accent_color=ACCENT3)

    # =========================================================================
    # SLIDE 4: Sticking Mechanic
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 2, "The Sticking Mechanic")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5),
                 "After ANY turn, all players can attempt to stick cards matching the discard top",
                 font_size=15, color=LIGHT)

    add_card(slide, Inches(0.5), Inches(2.2), Inches(4.2), Inches(1.5),
             "✓ SUCCESSFUL STICK", [
                 "Card rank matches discard top",
                 "Card is removed from hand permanently",
                 "Hand shrinks → lower final score",
                 "Huge advantage if used well",
             ], accent_color=ACCENT3)

    add_card(slide, Inches(5.0), Inches(2.2), Inches(4.5), Inches(1.5),
             "✗ FAILED STICK", [
                 "Card rank does NOT match",
                 "Penalty: draw a random card from deck",
                 "Hand GROWS → higher final score",
                 "High risk if guessing on unknown cards",
             ], accent_color=ACCENT2)

    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(1),
                 "Key insight: Base agents never stick. Smart sticks known matches only.\n"
                 "Card-counting agents add probabilistic sticking, ending rounds with 3.1 cards on average vs 4.0.",
                 font_size=13, color=ACCENT, bold=False)

    # =========================================================================
    # SLIDE 5: POMDP Formalization
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 3, "Formal Problem: A POMDP")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5),
                 "Cambio is a Partially Observable Markov Decision Process (S, A, T, R, Ω, O)",
                 font_size=14, color=LIGHT)

    cards_data = [
        ("State Space S", "Full deck ordering, all hands,\ndiscard pile, turn — mostly HIDDEN", ACCENT),
        ("Action Space A", "Draw, swap, discard, use power,\nstick, call Cambio", ACCENT),
        ("Observations Ω", "Initial peek, ability reveals,\ndiscard pile, opponent actions", ACCENT2),
    ]
    for i, (title, body, color) in enumerate(cards_data):
        add_card(slide, Inches(0.5 + i * 3.1), Inches(2.2), Inches(2.9), Inches(1.5),
                 title, body.split('\n'), accent_color=color)

    add_text_box(slide, Inches(0.5), Inches(4.0), Inches(9), Inches(1),
                 "54! possible deck orderings. Exact POMDP solving is intractable.\n"
                 "Our approach: deterministic tracking of all observed cards + expected value over the rest.",
                 font_size=13, color=DIM)

    # =========================================================================
    # SLIDE 6: System Architecture
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 4, "System Architecture")

    # --- Architecture diagram with arrows ---
    # Top row: Game Engine (center)
    add_card(slide, Inches(3.2), Inches(1.4), Inches(3.6), Inches(1.2),
             "GAME ENGINE", [
                 "Rules, deck/discard, power execution",
                 "Cambio bonus/penalty, information hiding",
             ], accent_color=ACCENT)

    # Middle row: Agents (left) and Web UI (right), both feed into engine
    add_card(slide, Inches(0.3), Inches(3.0), Inches(3.0), Inches(1.6),
             "AI AGENTS (agents/)", [
                 "Base → Smart → V1 → V2",
                 "Shared interface: choose_draw,",
                 "choose_action, call_cambio, etc.",
                 "V2: 9 features + 2 tuned presets",
             ], accent_color=ACCENT3)

    add_card(slide, Inches(6.7), Inches(3.0), Inches(3.0), Inches(1.6),
             "WEB INTERFACE (web/)", [
                 "Flask + SocketIO real-time UI",
                 "Human-vs-AI play via browser",
                 "Event log, card animations",
                 "Same info constraints as bots",
             ], accent_color=RGBColor(0xDD, 0xA0, 0xFF))

    # Bottom: Tournament + Benchmark
    add_card(slide, Inches(3.2), Inches(3.0), Inches(3.2), Inches(1.6),
             "TOURNAMENT SYSTEM", [
                 "Multi-round matches, first to 100 loses",
                 "21 matchup configs × 100-200 matches",
                 "Ablation study + grid search",
                 "Randomized seating, per-round stats",
             ], accent_color=ACCENT2)

    # Arrows: agents → engine, web → engine, engine → tournament
    def add_arrow(slide, x1, y1, x2, y2, color=DIM):
        connector = slide.shapes.add_connector(
            1,  # straight connector
            Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = color
        connector.line.width = Pt(2)

    add_arrow(slide, 1.8, 3.0, 3.4, 2.6, ACCENT3)   # agents → engine
    add_arrow(slide, 8.2, 3.0, 6.6, 2.6, RGBColor(0xDD, 0xA0, 0xFF))  # web → engine
    add_arrow(slide, 5.0, 2.6, 5.0, 3.0, ACCENT2)    # engine → tournament

    # =========================================================================
    # SLIDE 7: Agent Progression
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "Agent Progression")

    add_text_box(slide, Inches(0.5), Inches(1.4), Inches(9), Inches(0.4),
                 "Each agent builds on the previous — from simple heuristics to belief-state reasoning",
                 font_size=14, color=DIM)

    agents = [
        ("Base Agent", "Heuristic draw/swap\nUses all powers\nNever calls Cambio\nNever sticks", "BASELINE", DIM),
        ("Smart Agent", "Threshold-based draw\nKnown-card swaps only\nEstimated Cambio timing\nSticks known matches", "HEURISTIC", RGBColor(0xFF, 0xAA, 0x33)),
        ("Bayesian V1", "Tracks all 54 cards\nEV-based decisions\nAdaptive Cambio\nKnown-card sticking", "CARD COUNTING", ACCENT),
        ("Bayesian V2", "+ Opponent modeling\n+ Disruption swaps\n+ Probabilistic sticking\n+ Final-round mode\n+ Deck awareness", "OPPONENT AWARE", ACCENT3),
    ]
    for i, (name, body, label, color) in enumerate(agents):
        left = Inches(0.3 + i * 2.4)
        add_card(slide, left, Inches(1.9), Inches(2.2), Inches(2.5),
                 name, body.split('\n'), accent_color=color)
        add_text_box(slide, left, Inches(4.45), Inches(2.2), Inches(0.4),
                     label, font_size=9, color=color, bold=True, alignment=PP_ALIGN.CENTER)

        if i < 3:
            add_text_box(slide, left + Inches(2.05), Inches(2.8), Inches(0.5), Inches(0.5),
                         "→", font_size=24, color=DIM, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 8: Smart Agent Deep Dive
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "Smart Agent — Rule-Based Strategy")

    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(1.2),
             "DRAW DECISION", [
                 "Discard pile value < 4 → take from discard",
                 "Otherwise → draw from deck",
             ])

    add_card(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(1.2),
             "SWAP DECISION", [
                 "Compare drawn card to KNOWN positions only",
                 "Swap into highest-value known position",
                 "Cannot reason about unknown positions",
             ], accent_color=ACCENT2)

    add_card(slide, Inches(0.5), Inches(2.9), Inches(4.3), Inches(1.6),
             "CAMBIO TIMING", [
                 "Estimates own score: known cards + 5/unknown",
                 "Estimates opponent: known cards + 6/unknown",
                 "Calls when: score < 10 AND margin > 4",
                 "Also calls if all cards known AND total < 8",
             ])

    add_card(slide, Inches(5.0), Inches(2.9), Inches(4.5), Inches(1.6),
             "KEY WEAKNESS", [
                 "No estimate for unknown positions",
                 "Uses fixed average (5) for unknowns",
                 "Sometimes passes up beneficial swaps",
                 "May call Cambio hiding a 10-point card",
             ], accent_color=ACCENT2)

    add_text_box(slide, Inches(0.5), Inches(4.7), Inches(9), Inches(0.5),
                 "Wins 63% vs Base (meaningful), but gets crushed by card-counting agents",
                 font_size=12, color=DIM)

    # =========================================================================
    # SLIDE 9: V1 - Card Tracking
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "Bayesian V1 — Tracking All 54 Cards")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5),
                 "Core idea: maintain a CardTracker that knows WHERE every observed card is",
                 font_size=15, color=ACCENT, bold=True)

    add_card(slide, Inches(0.5), Inches(2.1), Inches(2.8), Inches(2.0),
             "ACCOUNTED CARDS", [
                 "Discard pile (synced each turn)",
                 "Own hand (known positions)",
                 "Opponent hands (from peeks",
                 "  and observed discard draws)",
                 "",
                 "All tracked as (rank, suit) tuples",
             ], accent_color=ACCENT3)

    add_card(slide, Inches(3.5), Inches(2.1), Inches(2.8), Inches(2.0),
             "UNACCOUNTED POOL", [
                 "U = full 54-card deck",
                 "  − discard pile",
                 "  − known own cards",
                 "  − known opponent cards",
                 "",
                 "Any unknown position's EV =",
                 "  mean value of U",
             ], accent_color=ACCENT2)

    add_card(slide, Inches(6.5), Inches(2.1), Inches(3.0), Inches(2.0),
             "THE KEY FORMULA", [
                 "",
                 "E[unknown] = (1/|U|) Σ value(c)",
                 "             for all c ∈ U",
                 "",
                 "As the game progresses,",
                 "|U| shrinks → estimates get",
                 "MORE precise over time",
             ], accent_color=RGBColor(0xDD, 0xA0, 0xFF))

    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.8),
                 "V1 also adds: EV-based swap targeting (swap into highest-EV position),\n"
                 "information bonus (+1 for placing known-low card into unknown slot),\n"
                 "adaptive Cambio timing (threshold 10, EV dominance margin 8)",
                 font_size=12, color=DIM)

    # =========================================================================
    # SLIDE 10: V1 Decision Strategy
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "V1 — Decision Strategy")

    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(1.4),
             "DRAW DECISION", [
                 "Take from discard only when:",
                 "  • Card is Joker (0) or Red King (-1), OR",
                 "  • Improvement ≥ 1 over worst expected position",
                 "Threshold of 1 accounts for lost deck-draw optionality",
             ])

    add_card(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(1.4),
             "SWAP DECISION", [
                 "For each position, compute:",
                 "  improvement = E[position] − drawn_value",
                 "  + info_bonus (1 if drawn ≤ 3 into unknown)",
                 "Swap into position with max improvement (if > 0)",
             ], accent_color=ACCENT)

    add_card(slide, Inches(0.5), Inches(3.1), Inches(4.3), Inches(1.4),
             "PEEK TARGETING", [
                 "7/8: peek first unknown own position",
                 "9/10: peek lowest-EV opponent (most",
                 "  threatening), random unknown position",
                 "Black King: peek opponent, swap only if better",
             ], accent_color=RGBColor(0xDD, 0xA0, 0xFF))

    add_card(slide, Inches(5.0), Inches(3.1), Inches(4.5), Inches(1.4),
             "CAMBIO TIMING (two paths)", [
                 "Path A — High confidence:",
                 "  Know ≥ 3 of 4 cards, score < 10, margin > 4",
                 "Path B — EV dominance:",
                 "  Expected score ≥ 8 pts below ALL opponents",
             ], accent_color=ACCENT3)

    # =========================================================================
    # SLIDE 11: V2 Improvements Overview
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "Bayesian V2 — Nine Improvements")

    improvements = [
        ("1. Opponent Self-Knowledge Tracking",
         "Track which hand positions each opponent likely knows about.\n"
         "After deal: everyone knows positions 0,1. Draw+swap → gain knowledge.\n"
         "Blind/king swap → lose knowledge."),
        ("2. Disruption-Weighted Swap Scoring",
         "score(p) = -value(p) + disruption_bonus(p)\n"
         "Bonus up to 3 (scaled to 1 in 1v1) when opponent knows position p.\n"
         "Targets exact set of positions opponent tracks, forcing information loss."),
        ("3. Third-Party Swaps",
         "When own hand is strong (worst ≤ 5) AND we're winning, use J/Q to\n"
         "swap TWO OPPONENTS' known-low positions. Only disrupts when ahead;\n"
         "skips sabotage when behind to focus on self-improvement."),
    ]

    for i, (title, body) in enumerate(improvements):
        add_card(slide, Inches(0.5), Inches(1.45 + i * 1.3), Inches(9.0), Inches(1.15),
                 title, body.split('\n'),
                 accent_color=[ACCENT, ACCENT2, ACCENT3][i])

    # =========================================================================
    # SLIDE 12: V2 Improvements (cont)
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "V2 Improvements (continued)")

    improvements2 = [
        ("4. Probabilistic Stick Play",
         "V1 only sticks known matches. V2 sticks UNKNOWN positions when:\n"
         "EV(stick) = P(match) × card_value − (1−P(match)) × E[penalty] > 2.0\n"
         "P(match) = count of matching rank in unaccounted pool / |U|"),
        ("5. Improved Cambio Timing",
         "Aggressive threshold: 6 (vs V1's 8) when all cards known.\n"
         "Cambio margin: 3 (vs V1's 4). Knowledge gap: 2 (vs V1's 1).\n"
         "Preemptive calling tested but DISABLED — ablation showed it hurts."),
        ("6. Smarter Peek Targeting (9/10)",
         "Target lowest-EV opponent (most dangerous), then peek UNMANAGED\n"
         "unknown positions first (not in opponent's self-knowledge set).\n"
         "Unmanaged positions more likely hold high-value unswapped cards."),
    ]

    for i, (title, body) in enumerate(improvements2):
        add_card(slide, Inches(0.5), Inches(1.45 + i * 1.3), Inches(9.0), Inches(1.15),
                 title, body.split('\n'),
                 accent_color=[RGBColor(0xDD, 0xA0, 0xFF), ACCENT, ACCENT3][i])

    # =========================================================================
    # SLIDE 12b: V2 Advanced Features (7-9)
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 5, "V2 Advanced Features")

    improvements3 = [
        ("7. Final-Round Mode",
         "When someone calls Cambio, V2 shifts to end-game strategy:\n"
         "• Take ANY discard improvement (threshold 0 instead of 1)\n"
         "• Suppress probabilistic sticks (penalty card = direct score hit)\n"
         "• More aggressive swaps into unknown positions"),
        ("8. Opponent Action Inference",
         "Track opponent behavior to refine EV estimates:\n"
         "• Opponent discards without swapping → all their cards ≤ discard value\n"
         "• Opponent peeks own card (7/8) and keeps → that position is likely low\n"
         "• Uses Bayesian priors to lower opponent score estimates"),
        ("9. Deck Size Awareness",
         "Monitor len(deck) to adjust strategy in late-game:\n"
         "• Deck ≤ 5 cards → prefer known discard over risky deck draw\n"
         "• Nearly empty deck + ahead of all opponents → call Cambio early\n"
         "  (reshuffle increases variance; lock in the lead)"),
    ]

    for i, (title, body) in enumerate(improvements3):
        add_card(slide, Inches(0.5), Inches(1.45 + i * 1.3), Inches(9.0), Inches(1.15),
                 title, body.split('\n'),
                 accent_color=[ACCENT, ACCENT2, ACCENT3][i])

    # =========================================================================
    # SLIDE 13: Experiment Design
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 6, "Experiment Design")

    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(2.0),
             "TOURNAMENT SYSTEM", [
                 "17 matchup configurations",
                 "100 multi-round matches each",
                 "First to 100 cumulative points loses",
                 "Randomized seating order each round",
                 "~5-9 rounds per match on average",
                 "",
                 "Configs span: 1v1, 3-player, 6-player",
             ])

    add_card(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(2.0),
             "METRICS", [
                 "Match Win Rate — fraction of matches won",
                 "Per-Round Avg Score — consistency measure",
                 "Per-Round Avg Hand Size — sticking success",
                 "Cambio Caller Win Rate — timing accuracy",
                 "Discard Draw Rate — draw strategy",
                 "Swap Rate — swap aggressiveness",
             ], accent_color=ACCENT2)

    add_text_box(slide, Inches(0.5), Inches(3.8), Inches(9), Inches(1.2),
                 "Key matchups:\n"
                 "  • BayesV2 vs Smart (1v1) — primary evaluation\n"
                 "  • 2×Smart + BayesV2 (3p) — multiplayer scaling\n"
                 "  • 5×Smart + BayesV2 (6p) — large table scaling\n"
                 "  • BayesV2 vs BayesV1 (1v1) — measures V2 improvement\n"
                 "  • Smart vs Base (1v1 control) — validates Smart plays meaningfully",
                 font_size=12, color=DIM)

    # =========================================================================
    # SLIDE 14: Main Results
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 7, "Results: Card-Counting vs Heuristic")

    add_stat_box(slide, Inches(0.5), Inches(1.5), "92%", "1v1 WIN RATE", ACCENT)
    add_stat_box(slide, Inches(2.5), Inches(1.5), "50%", "3-PLAYER WIN RATE", ACCENT)
    add_stat_box(slide, Inches(4.5), Inches(1.5), "42%", "4-PLAYER WIN RATE", ACCENT)
    add_stat_box(slide, Inches(6.8), Inches(1.5), "56%", "V2 vs V1 (1v1)", ACCENT3)
    add_stat_box(slide, Inches(8.3), Inches(1.5), "52%", "V2 vs OLD V2", DIM)

    add_text_box(slide, Inches(0.5), Inches(2.6), Inches(9), Inches(0.4),
                 "BayesV2 with optimized presets (duel/multi) — 200 matches per config, mixed opponents",
                 font_size=11, color=DIM)

    # Results table
    add_card(slide, Inches(0.5), Inches(3.0), Inches(9.0), Inches(2.2),
             "MATCH WIN RATES (200 matches, optimized presets)", [
                 "",
                 "1v1:  V2-Duel  92% vs Smart 8%",
                 "3p:   V2-Multi 50% vs V1 37%, Smart 12%",
                 "4p:   V2-Multi 42% vs V1 36%, Smart 23%, Base 0%",
                 "",
                 "Control: Smart 90% vs Base 10% — card tracking is the core structural advantage",
                 "Fair-game baselines: 50% (1v1), 33% (3p), 25% (4p)",
             ])

    # =========================================================================
    # SLIDE 15: V1 vs V2 Results
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 7, "Results: V1 vs V2")

    add_stat_box(slide, Inches(0.5), Inches(1.5), "56%", "V2-DUEL (1v1)", ACCENT3)
    add_stat_box(slide, Inches(2.5), Inches(1.5), "50%", "V2-MULTI (3p)", ACCENT3)
    add_stat_box(slide, Inches(4.5), Inches(1.5), "32%", "V2-MULTI (4p)", ACCENT)
    add_stat_box(slide, Inches(6.5), Inches(1.5), "52%", "V2 vs DEFAULT", DIM)

    add_card(slide, Inches(0.5), Inches(2.7), Inches(4.3), Inches(2.2),
             "V2-DUEL vs V1 (1v1)", [
                 "56% win rate (V2-Duel) vs 44% (V1)",
                 "52% win rate vs V2-Default",
                 "",
                 "Advanced features + tuned parameters",
                 "give consistent edge even against",
                 "another card-counting agent",
             ], accent_color=ACCENT3)

    add_card(slide, Inches(5.0), Inches(2.7), Inches(4.5), Inches(2.2),
             "V2-MULTI IN MULTIPLAYER", [
                 "3p: V2-Multi 50% vs V1 37%, Smart 12%",
                 "4p: V2-Multi 32% vs V2-Def 24%,",
                 "    V1 26%, Smart 18%",
                 "",
                 "Wrong preset hurts: duel in 3p → 35%",
                 "vs multi in 3p → 56%. Adaptation to",
                 "table size is a real advantage.",
             ], accent_color=ACCENT3)

    # =========================================================================
    # SLIDE 16: Gameplay Stats
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 7, "Per-Round Gameplay Analysis")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 "BayesV2 vs Smart (1v1) — 729 rounds across 100 matches",
                 font_size=13, color=DIM)

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.3), Inches(2.8),
             "WHERE V2 WINS", [
                 "Avg cards at round end: 3.1 vs 4.0",
                 "  → V2 sticks probabilistically; Smart only known",
                 "",
                 "Cambio call rate: 84% vs 10%",
                 "  → V2 controls when rounds end",
                 "",
                 "Cambio caller win rate: 74% vs 53%",
                 "  → V2's timing is more accurate",
                 "",
                 "Score per round: 5.3 vs 14.7",
                 "  → ~10 point advantage EVERY round",
             ], accent_color=ACCENT3)

    add_card(slide, Inches(5.0), Inches(2.0), Inches(4.5), Inches(2.8),
             "WHERE THEY'RE SIMILAR", [
                 "Discard draw rate: 29% vs 31%",
                 "Swap rate: 35% vs 36%",
                 "Power usage: 3.6 vs 3.4 per round",
                 "",
                 "The advantage isn't from acting",
                 "MORE — it's from acting SMARTER.",
                 "",
                 "Same number of draws and swaps,",
                 "but V2 targets the RIGHT positions",
                 "and knows WHEN to call Cambio.",
             ], accent_color=ACCENT2)

    # =========================================================================
    # SLIDE 17: Human Evaluation
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 7, "Human Evaluation")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.6),
                 "Each team member played 25 games against BayesV2 via the web interface (75 total)",
                 font_size=14, color=LIGHT)

    add_stat_box(slide, Inches(0.5), Inches(2.3), "56%", "HANNAH", ACCENT)
    add_stat_box(slide, Inches(2.5), Inches(2.3), "52%", "NAS", ACCENT)
    add_stat_box(slide, Inches(4.5), Inches(2.3), "44%", "JUAN", ACCENT2)
    add_stat_box(slide, Inches(7.0), Inches(2.3), "7.7", "AGENT AVG SCORE", DIM)

    add_card(slide, Inches(0.5), Inches(3.4), Inches(4.3), Inches(1.6),
             "NOT STATISTICALLY SIGNIFICANT", [
                 "n=25 per player is too small",
                 "Binomial test: all p > 0.05",
                 "Would need ~200+ games each for",
                 "significance at this effect size",
                 "Treated as qualitative observation",
             ], accent_color=DIM)

    add_card(slide, Inches(5.0), Inches(3.4), Inches(4.5), Inches(1.6),
             "QUALITATIVE OBSERVATIONS", [
                 "Agent's superpower: perfect memory",
                 "  → never forgets a peeked card",
                 "Human's superpower: intuition",
                 "  → better timing, aggressive play",
                 "Result: roughly even → human-level",
             ], accent_color=ACCENT3)

    # =========================================================================
    # SLIDE 18: Web UI
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 8, "Web Interface")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5),
                 "Flask + SocketIO real-time multiplayer interface",
                 font_size=15, color=LIGHT)

    add_card(slide, Inches(0.5), Inches(2.1), Inches(4.3), Inches(2.5),
             "FEATURES", [
                 "• Lobby system with game creation",
                 "• Real-time card animations",
                 "• Clickable action buttons",
                 "  (Draw from Deck / Take Discard / Cambio)",
                 "• Full event log of all game actions",
                 "• Visible discard pile + own cards",
                 "• AI opponents play automatically",
                 "• Same information constraints as bots",
             ])

    add_card(slide, Inches(5.0), Inches(2.1), Inches(4.5), Inches(2.5),
             "ARCHITECTURE", [
                 "web/app.py — Flask + SocketIO server",
                 "web/game_manager.py — game lifecycle",
                 "web/human_agent.py — bridges browser",
                 "  input to game engine's Player interface",
                 "",
                 "Human turns block (async wait for",
                 "browser input); AI turns resolve",
                 "instantly and broadcast via WebSocket",
             ], accent_color=ACCENT2)

    add_text_box(slide, Inches(0.5), Inches(4.8), Inches(9), Inches(0.5),
                 "[ LIVE DEMO ]",
                 font_size=20, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 19: Ablation Study
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 9, "Ablation Study: V2 Feature Contributions")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 "Disable each V2 feature one at a time; measure win rate vs Smart (100 matches each)",
                 font_size=13, color=DIM)

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.3), Inches(3.0),
             "1v1 vs Smart (baseline 90%)", [
                 "",
                 "Feature Off         WR   Δ",
                 "───────────      ───── ────",
                 "Disruption         80%  -10%",
                 "Prob. Stick        81%   -9%",
                 "Opp Inference      87%   -3%",
                 "Deck Awareness     86%   -4%",
                 "Final-Round        89%   -1%",
                 "Aggr. Cambio       88%   -2%",
                 "Smart Peek         88%   -2%",
                 "3rd-Party Swaps    89%   -1%",
                 "Preemptive Cmb     90%   +0%",
             ])

    add_card(slide, Inches(5.0), Inches(2.0), Inches(4.5), Inches(3.0),
             "3-PLAYER vs Smart (baseline 72%)", [
                 "",
                 "Feature Off         WR    Δ",
                 "───────────      ─────  ────",
                 "Prob. Stick        59%  -13%",
                 "Smart Peek         60%  -12%",
                 "Disruption         62%  -10%",
                 "Opp Inference      62%  -10%",
                 "Preemptive Cmb     62%  -10%",
                 "3rd-Party Swaps    62%  -10%",
                 "Final-Round        65%   -7%",
                 "Deck Awareness     65%   -7%",
                 "Aggr. Cambio       76%   +4%",
             ], accent_color=ACCENT3)

    add_text_box(slide, Inches(0.5), Inches(5.1), Inches(9), Inches(0.4),
                 "Top contributors: Disruption scoring (-10%), Prob. stick (-9%/-13%), Opponent inference (-10% in 3p)",
                 font_size=11, color=ACCENT)

    # =========================================================================
    # SLIDE 20: Grid Search — Parameter Tuning
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 9, "Parameter Grid Search")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 "11 parameters swept independently × 3 table sizes × 100 matches — 222 runs parallelized across 8 cores",
                 font_size=12, color=DIM)

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.3), Inches(3.2),
             "KEY PARAMETERS EXPLAINED", [
                 "disruption_bonus — bonus for swapping",
                 "  positions the opponent tracks (+info loss)",
                 "cambio_threshold — max score to call Cambio",
                 "ev_dominance_margin — how far ahead before",
                 "  calling Cambio on expected value alone",
                 "stick_ev_threshold — min EV to attempt a",
                 "  probabilistic stick on an unknown card",
                 "good_hand_threshold — max worst-card value",
                 "  to switch from self-improvement to disruption",
                 "jq_swap_improvement — min point gain to use",
                 "  J/Q on a known opponent target",
                 "small_deck_threshold — deck size below which",
                 "  we prefer safe discard draws over deck risk",
             ])

    add_card(slide, Inches(5.0), Inches(2.0), Inches(4.5), Inches(3.2),
             "BIGGEST IMPROVEMENTS vs DEFAULT", [
                 "",
                 "Parameter             1v1   4p",
                 "─────────────       ───── ─────",
                 "cambio_threshold     +12%  +14%",
                 "small_deck_thresh    +10%  +14%",
                 "disruption_bonus      +9%  +14%",
                 "stick_ev_thresh       +8%  +18%",
                 "cambio_margin         +5%   +9%",
                 "ev_dominance          +7%   +4%",
                 "jq_swap_improve       +6%  +12%",
                 "",
                 "Mixed opponents: Smart + V1 + Base",
             ], accent_color=ACCENT3)

    # =========================================================================
    # SLIDE 21: Duel vs Multi Presets
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 9, "Adaptive Playstyles: Duel vs Multi")

    add_text_box(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4),
                 "Optimal parameters differ sharply by table size — so V2 adapts automatically",
                 font_size=13, color=LIGHT)

    add_card(slide, Inches(0.5), Inches(2.0), Inches(4.3), Inches(2.0),
             "DUEL PRESET (1v1)", [
                 "High disruption (4) — one opponent to target",
                 "Aggressive Cambio (threshold 9, margin 4)",
                 "Loose sticking (EV > 1.5) — few penalty sources",
                 "Low swap bar (improvement ≥ 1 for J/Q)",
                 "Tight good-hand threshold (2) — disrupt early",
                 "",
                 "→ 92% vs Smart,  56% vs V1",
             ], accent_color=ACCENT)

    add_card(slide, Inches(5.0), Inches(2.0), Inches(4.5), Inches(2.0),
             "MULTI PRESET (3-4 player)", [
                 "Near-zero disruption (0.5) — helps other opps",
                 "Conservative Cambio (threshold 8, margin 5)",
                 "Tighter sticking (EV > 2.0) — more penalty risk",
                 "High swap bar (improvement ≥ 4 for J/Q)",
                 "Wide good-hand threshold (8) — self-improve first",
                 "",
                 "→ 50% in 3p,  42% in 4p (vs mixed field)",
             ], accent_color=ACCENT3)

    add_card(slide, Inches(0.5), Inches(4.2), Inches(9.0), Inches(1.0),
             "WHY DISRUPTION FLIPS", [
                 "In 1v1, swapping the opponent's known position destroys THEIR info and only helps YOU.",
                 "In 4p, the same swap helps the OTHER 2 opponents — they benefit from your opponent's confusion.",
                 "Grid search found disruption_bonus=4 in 1v1 but 0.5 in 4p: a complete reversal.",
             ], accent_color=ACCENT2)

    # =========================================================================
    # SLIDE: Key Findings
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 10, "Key Findings")

    findings = [
        ("1", "Classical probabilistic reasoning dominates heuristics",
         "92% win rate (tuned) without any training, neural networks, or search trees. Just tracking + expected value + grid search.",
         ACCENT),
        ("2", "Optimal strategy depends sharply on table size",
         "Disruption bonus flips from 4 (1v1) to 0.5 (4p). Using the wrong preset costs ~20pp. The agent must adapt.",
         ACCENT3),
        ("3", "Parameter tuning unlocks massive gains",
         "11-parameter grid search improved 1v1 from 82% → 92% (+10pp). Defaults left significant performance on the table.",
         ACCENT2),
        ("4", "Sticking is the single biggest structural advantage",
         "Agents that stick end with 3.1 cards vs 4.0 — a ~25% reduction in hand size and final score.",
         RGBColor(0xDD, 0xA0, 0xFF)),
    ]

    for i, (num, title, body, color) in enumerate(findings):
        y = Inches(1.45 + i * 0.95)
        add_text_box(slide, Inches(0.5), y, Inches(0.5), Inches(0.4),
                     num, font_size=24, color=color, bold=True)
        add_text_box(slide, Inches(1.1), y, Inches(8), Inches(0.3),
                     title, font_size=15, color=WHITE, bold=True)
        add_text_box(slide, Inches(1.1), y + Inches(0.3), Inches(8), Inches(0.4),
                     body, font_size=11, color=DIM)

    # =========================================================================
    # SLIDE 20: Limitations & Future Work
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 10, "Limitations & Future Work")

    add_card(slide, Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.2),
             "LIMITATIONS", [
                 "• No individual opponent modeling",
                 "  (assumes all opponents play similarly)",
                 "",
                 "• Inference priors are coarse-grained",
                 "  (upper bounds, not per-card posteriors)",
                 "",
                 "• J/Q can't swap same-player positions",
                 "",
                 "• Human evaluation too small for",
                 "  statistical significance (n=25)",
                 "",
                 "• Grid search is one-at-a-time (not joint)",
             ], accent_color=ACCENT2)

    add_card(slide, Inches(5.0), Inches(1.5), Inches(4.5), Inches(3.2),
             "FUTURE WORK", [
                 "• Reinforcement learning agents",
                 "  Can RL discover non-obvious strategies?",
                 "",
                 "• Per-opponent behavior models",
                 "  Adapt swap targeting to individual styles",
                 "",
                 "• Larger human studies",
                 "  100+ games/player, naive participants",
                 "",
                 "• Joint parameter optimization (Bayesian opt)",
                 "  Current search is univariate; interactions",
                 "  between parameters remain unexplored",
             ], accent_color=ACCENT3)

    # =========================================================================
    # SLIDE 21: Summary
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)

    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(9), Inches(0.8),
                 "What We Built",
                 font_size=32, color=WHITE, bold=True)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(2), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    items = [
        "A complete game engine faithfully implementing all Cambio rules and powers",
        "4 progressively sophisticated AI agents (Base → Smart → V1 → V2)",
        "A CardTracker that monitors all 54 cards across every game event",
        "9-feature V2 with adaptive duel/multi presets tuned via grid search",
        "11-parameter grid search (222 configs × 100 matches, parallelized)",
        "A tournament system running 21 configurations × 100-200 matches",
        "A real-time web interface for human-vs-AI play",
        "92% 1v1 win rate — from 82% default to 92% after optimization",
    ]
    add_bullet_list(slide, Inches(0.7), Inches(1.3), Inches(8.5), Inches(3.0),
                    items, font_size=14, color=LIGHT)

    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(9), Inches(0.8),
                 "Classical AI techniques — no training data, no neural networks —\n"
                 "can achieve human-level play in imperfect-information card games.",
                 font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

    # =========================================================================
    # SLIDE 22: References
    # =========================================================================
    slide = prs.slides.add_slide(blank_layout)
    add_bg(slide)
    add_section_header(slide, 11, "References")

    refs = [
        "N. Bard et al., \"The Hanabi Challenge,\" Artificial Intelligence, vol. 280, 2020.",
        "C. Browne et al., \"A Survey of MCTS Methods,\" IEEE Trans. CI & AI in Games, 2012.",
        "R. Sutton and A. Barto, Reinforcement Learning: An Introduction, 2nd ed., MIT Press, 2018.",
        "A. Raffin et al., \"Stable-Baselines3,\" JMLR, vol. 22, no. 268, 2021.",
        "M. Ginsberg, \"GIB: Expert-Level Bridge,\" IJCAI, 1999.",
        "N. Brown and T. Sandholm, \"Superhuman AI for Multiplayer Poker,\" Science, 2019.",
        "M. Moravčík et al., \"DeepStack,\" Science, 2017.",
        "M. Zinkevich et al., \"Regret Minimization in Games,\" NeurIPS, 2007.",
        "L. Kaelbling et al., \"Planning and Acting in POMDPs,\" AI, vol. 101, 1998.",
        "Cambio Card Game Rules, cambiocardgame.com.",
        "A. Wu, \"How to Play Cambio,\" wikiHow, Sept. 2025.",
        "Nyhead, \"CambioMinimax,\" GitHub.",
    ]

    add_bullet_list(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(3.5),
                    refs, font_size=10, color=DIM)

    # =========================================================================
    # Save
    # =========================================================================
    out = 'Cambio_AI_Final_Presentation.pptx'
    prs.save(out)
    print(f'Saved {out} ({len(prs.slides)} slides)')


if __name__ == '__main__':
    build()
